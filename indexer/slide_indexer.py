"""
Slide indexer - Creates slide_index.jsonl from sessions.

This module handles:
- Downloading PPTX files
- Extracting text content from slides
- Creating JSONL index files for Azure AI Search
"""

import asyncio
import gc
import json
import logging
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Optional

import aiohttp

from .models import SessionInfo, SlideRecord, HTTP_HEADERS, IGNORE_SESSION_CODES

logger = logging.getLogger(__name__)

# Thread pool for blocking PPTX operations
_executor = ThreadPoolExecutor(max_workers=10)


def extract_text_from_slide(slide) -> str:
    """Extract all text content from a PPTX slide object."""
    text_runs = []
    
    try:
        # Title
        if slide.shapes.title and slide.shapes.title.text:
            text_runs.append(slide.shapes.title.text.strip())
        
        # Shapes and Tables
        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                continue
            
            # Text Boxes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_runs.append(text)
            
            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text_frame.text.strip()
                        if text:
                            text_runs.append(text)
                            
    except Exception:
        pass  # Skip malformed shapes
    
    return "\n".join(text_runs)


def parse_pptx_file(filepath: Path) -> list[tuple[int, str]]:
    """
    Parse a PPTX file and extract slide content.
    
    Args:
        filepath: Path to the PPTX file
    
    Returns:
        List of (slide_number, content) tuples
    """
    from pptx import Presentation
    
    slides_data = []
    prs = None
    
    try:
        prs = Presentation(filepath)
        for slide_idx, slide in enumerate(prs.slides):
            content = extract_text_from_slide(slide)
            if content.strip():
                slides_data.append((slide_idx + 1, content))
        return slides_data
    finally:
        if prs:
            del prs
        gc.collect()


async def download_pptx(
    http_session: aiohttp.ClientSession,
    url: str,
    dest_path: Path
) -> bool:
    """Download a PPTX file from URL."""
    try:
        async with http_session.get(url, headers=HTTP_HEADERS) as resp:
            if resp.status == 200:
                content = await resp.read()
                dest_path.parent.mkdir(parents=True, exist_ok=True)
                dest_path.write_bytes(content)
                return True
            else:
                logger.warning(f"Download failed ({resp.status}): {url}")
    except Exception as e:
        logger.error(f"Download error for {url}: {e}")
    return False


async def process_session(
    http_session: aiohttp.ClientSession,
    session: SessionInfo,
    ppts_dir: Path,
    semaphore: asyncio.Semaphore,
) -> list[SlideRecord]:
    """
    Process a single session - download PPTX and extract slides.
    
    For Partner sessions: Creates a single record with description content
    (Partner Marketing Center requires authentication for PPTX downloads).
    
    For Build/Ignite sessions: Downloads PPTX and creates per-slide records.
    
    Args:
        http_session: aiohttp session
        session: Session to process
        ppts_dir: Directory to store PPTX files
        semaphore: Concurrency limiter
    
    Returns:
        List of SlideRecord objects
    """
    async with semaphore:
        if session.session_code in IGNORE_SESSION_CODES:
            logger.debug(f"Skipping ignored session: {session.session_code}")
            return []
        
        # Partner sessions: Can't download PPTX (requires auth)
        # Create a single record with title + description as content
        if session.event == 'Partner':
            content = session.title
            if session.description:
                content = f"{session.title}\n\n{session.description}"
            
            return [SlideRecord(
                slide_id=f"{session.session_code}_1",
                session_code=session.session_code,
                title=session.title,
                slide_number=1,
                content=content,
                event=session.event,
                session_url=session.session_url,
                ppt_url=session.ppt_url
            )]
        
        # Build/Ignite sessions: Download and parse PPTX
        filepath = ppts_dir / f"{session.session_code}.pptx"
        
        # Download if needed
        if not filepath.exists():
            if not session.ppt_url:
                return []
            
            logger.info(f"Downloading: {session.session_code}")
            success = await download_pptx(http_session, session.ppt_url, filepath)
            if not success:
                return []
        
        # Parse PPTX in thread pool
        try:
            loop = asyncio.get_event_loop()
            slides_data = await loop.run_in_executor(_executor, parse_pptx_file, filepath)
            
            records = []
            for slide_num, content in slides_data:
                records.append(SlideRecord(
                    slide_id=f"{session.session_code}_{slide_num}",
                    session_code=session.session_code,
                    title=session.title,
                    slide_number=slide_num,
                    content=content,
                    event=session.event,
                    session_url=session.session_url,
                    ppt_url=session.ppt_url
                ))
            
            logger.info(f"✓ {session.session_code}: {len(records)} slides")
            return records
            
        except Exception as e:
            logger.error(f"✗ Failed to parse {session.session_code}: {e}")
            
            # Delete corrupt files
            if "not a zip file" in str(e).lower() or "package not found" in str(e).lower():
                try:
                    filepath.unlink()
                    logger.info(f"Deleted corrupt file: {filepath}")
                except Exception:
                    pass
            
            return []


async def create_slide_index(
    sessions: list[SessionInfo],
    output_file: Path,
    ppts_dir: Path,
    download_ppts: bool = True,
    max_concurrent: int = 20,
) -> int:
    """
    Create slide_index.jsonl from sessions.
    
    Args:
        sessions: List of sessions to process
        output_file: Path to output JSONL file
        ppts_dir: Directory for PPTX files
        download_ppts: Whether to download and parse PPTX files
        max_concurrent: Max concurrent downloads
    
    Returns:
        Number of slide records written
    """
    output_file.parent.mkdir(parents=True, exist_ok=True)
    ppts_dir.mkdir(parents=True, exist_ok=True)
    
    total_records = 0
    
    if download_ppts:
        # Full extraction mode - download and parse PPTX files
        logger.info(f"Processing {len(sessions)} sessions (downloading PPTX files)...")
        
        timeout = aiohttp.ClientTimeout(total=120)
        connector = aiohttp.TCPConnector(limit=max_concurrent)
        semaphore = asyncio.Semaphore(max_concurrent)
        
        async with aiohttp.ClientSession(timeout=timeout, connector=connector) as http_session:
            tasks = [
                process_session(http_session, session, ppts_dir, semaphore)
                for session in sessions
            ]
            
            # Process in batches for progress reporting
            all_records = []
            batch_size = 20
            
            for i in range(0, len(tasks), batch_size):
                batch = tasks[i:i + batch_size]
                results = await asyncio.gather(*batch, return_exceptions=True)
                
                for result in results:
                    if isinstance(result, list):
                        all_records.extend(result)
                    elif isinstance(result, Exception):
                        logger.error(f"Task error: {result}")
                
                processed = min(i + batch_size, len(tasks))
                logger.info(f"Progress: {processed}/{len(tasks)} sessions")
            
            # Write all records to JSONL
            with open(output_file, 'w', encoding='utf-8') as f:
                for record in all_records:
                    f.write(json.dumps(record.to_dict(), ensure_ascii=False) + '\n')
            
            total_records = len(all_records)
    else:
        # Quick mode - create placeholder records without downloading
        logger.info(f"Creating placeholder index for {len(sessions)} sessions...")
        
        with open(output_file, 'w', encoding='utf-8') as f:
            for session in sessions:
                if session.session_code in IGNORE_SESSION_CODES:
                    continue
                
                record = SlideRecord(
                    slide_id=f"{session.session_code}_1",
                    session_code=session.session_code,
                    title=session.title,
                    slide_number=1,
                    content=session.title,
                    event=session.event,
                    session_url=session.session_url,
                    ppt_url=session.ppt_url
                )
                f.write(json.dumps(record.to_dict(), ensure_ascii=False) + '\n')
                total_records += 1
    
    logger.info(f"Written {total_records} records to {output_file}")
    return total_records


def load_sessions_from_jsonl(jsonl_path: Path) -> list[SessionInfo]:
    """
    Load sessions from an existing slide_index.jsonl file.
    
    Args:
        jsonl_path: Path to JSONL file
    
    Returns:
        List of unique SessionInfo objects
    """
    if not jsonl_path.exists():
        return []
    
    sessions = []
    seen_codes = set()
    
    with open(jsonl_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            
            try:
                data = json.loads(line)
                code = data.get('session_code')
                
                if code and code not in seen_codes:
                    seen_codes.add(code)
                    sessions.append(SessionInfo(
                        session_code=code,
                        title=data.get('title', ''),
                        event=data.get('event', ''),
                        session_id='',
                        session_url=data.get('session_url', ''),
                        ppt_url=data.get('ppt_url', '')
                    ))
            except json.JSONDecodeError:
                continue
    
    return sessions
