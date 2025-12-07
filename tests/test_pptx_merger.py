"""
Tests for PPTX Merger Service.
"""
import pytest
import zipfile
from pathlib import Path
from unittest.mock import patch, Mock
from pptx import Presentation
from services.pptx_merger import PPTXMerger

@pytest.fixture
def test_pptx_files(tmp_path):
    """Create two test PPTX files."""
    # Deck 1
    deck1_path = tmp_path / "deck1.pptx"
    prs1 = Presentation()
    slide_layout = prs1.slide_layouts[0] # Title Slide
    slide = prs1.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Deck 1 Slide 1"
    prs1.save(deck1_path)
    
    # Deck 2
    deck2_path = tmp_path / "deck2.pptx"
    prs2 = Presentation()
    slide_layout = prs2.slide_layouts[1] # Title and Content
    slide = prs2.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Deck 2 Slide 1"
    prs2.save(deck2_path)
    
    return deck1_path, deck2_path

def test_merge_pptx(test_pptx_files, tmp_path):
    """Test merging two PPTX files."""
    deck1, deck2 = test_pptx_files
    output_path = tmp_path / "merged.pptx"
    
    merger = PPTXMerger(output_path)
    merger.add_slide(deck1, 1)
    merger.add_slide(deck2, 1)
    merger.merge()
    
    assert output_path.exists()
    assert zipfile.is_zipfile(output_path)
    
    # Verify content with python-pptx
    prs = Presentation(output_path)
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "Deck 1 Slide 1"
    # Note: Text might be preserved, but styles are what we care about.
    # Since we copied XML, text should be there.
    assert prs.slides[1].shapes.title.text == "Deck 2 Slide 1"


def test_merge_single_slide(test_pptx_files, tmp_path):
    """Test merging a single slide."""
    deck1, _ = test_pptx_files
    output_path = tmp_path / "single.pptx"
    
    merger = PPTXMerger(output_path)
    merger.add_slide(deck1, 1)
    merger.merge()
    
    assert output_path.exists()
    prs = Presentation(output_path)
    assert len(prs.slides) == 1


def test_merge_empty_raises_error(tmp_path):
    """Test that merging with no slides raises an error."""
    output_path = tmp_path / "empty.pptx"
    
    merger = PPTXMerger(output_path)
    
    with pytest.raises(ValueError, match="No slides to merge"):
        merger.merge()


def test_merge_invalid_slide_number(test_pptx_files, tmp_path):
    """Test merging with invalid slide number."""
    deck1, _ = test_pptx_files
    output_path = tmp_path / "invalid.pptx"
    
    merger = PPTXMerger(output_path)
    merger.add_slide(deck1, 999)  # Non-existent slide
    
    # Should not crash but log a warning and produce a file with no slides
    # from that source (the merger skips invalid slides)
    merger.merge()
    
    # The output should still be created but may have fewer slides
    assert output_path.exists()


def test_pptx_merger_add_slide_from_same_source(test_pptx_files, tmp_path):
    """Test adding multiple slides from the same source."""
    deck1, _ = test_pptx_files
    
    # Create a deck with multiple slides
    prs = Presentation()
    for i in range(3):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i + 1}"
    multi_slide_path = tmp_path / "multi.pptx"
    prs.save(multi_slide_path)
    
    output_path = tmp_path / "merged_multi.pptx"
    merger = PPTXMerger(output_path)
    merger.add_slide(multi_slide_path, 1)
    merger.add_slide(multi_slide_path, 2)
    merger.add_slide(multi_slide_path, 3)
    merger.merge()
    
    assert output_path.exists()
    result_prs = Presentation(output_path)
    assert len(result_prs.slides) == 3
